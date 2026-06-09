#!/usr/bin/env python3
"""Generate PPTX from the canonical case spec."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
FONT = "Microsoft JhengHei"

COLORS = {
    "navy": RGBColor(25, 43, 79),
    "blue": RGBColor(49, 91, 163),
    "teal": RGBColor(34, 123, 157),
    "green": RGBColor(66, 122, 92),
    "gold": RGBColor(184, 138, 61),
    "red": RGBColor(180, 61, 54),
    "ink": RGBColor(40, 40, 40),
    "muted": RGBColor(95, 103, 115),
    "line": RGBColor(215, 220, 228),
    "bg": RGBColor(246, 248, 251),
    "white": RGBColor(255, 255, 255),
}


def load_spec(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as f:
        return json.load(f)


def table_rows(headers: list[str], rows: list[list[str]]) -> dict[str, Any]:
    return {"headers": headers, "rows": rows, "accent": "navy"}


class DeckBuilder:
    def __init__(self, spec: dict[str, Any]) -> None:
        self.spec = spec
        self.case = spec["case"]
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_W
        self.prs.slide_height = SLIDE_H
        self.layout = self.prs.slide_layouts[6]

    def color(self, name: str, fallback: str = "blue") -> RGBColor:
        return COLORS.get(name, COLORS[fallback])

    def add_bg(self, slide, color_name: str = "bg") -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.color(color_name, "bg")
        shape.line.fill.background()
        slide.shapes._spTree.remove(shape._element)
        slide.shapes._spTree.insert(2, shape._element)

    def add_header(self, slide, title: str, subtitle: str | None = None) -> None:
        title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(10.8), Inches(0.7))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = FONT
        r.font.size = Pt(24)
        r.font.bold = True
        r.font.color.rgb = COLORS["navy"]
        if subtitle:
            sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(0.96), Inches(11.3), Inches(0.45))
            p = sub_box.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = subtitle
            r.font.name = FONT
            r.font.size = Pt(10.5)
            r.font.color.rgb = COLORS["muted"]

    def add_footer(self, slide, text: str | None = None) -> None:
        footer_text = text or "Generated from canonical case spec"
        box = slide.shapes.add_textbox(Inches(0.55), Inches(7.0), Inches(12.1), Inches(0.25))
        p = box.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        r = p.add_run()
        r.text = footer_text
        r.font.name = FONT
        r.font.size = Pt(8.5)
        r.font.color.rgb = COLORS["muted"]

    def add_bullets(self, slide, items: list[str], x, y, w, h, *, size: int = 14, color: str = "ink") -> None:
        box = slide.shapes.add_textbox(x, y, w, h)
        tf = box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        first = True
        for item in items:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.space_after = Pt(6)
            r = p.add_run()
            r.text = item
            r.font.name = FONT
            r.font.size = Pt(size)
            r.font.color.rgb = self.color(color, "ink")

    def add_card(self, slide, title: str, bullets: list[str], accent: str, x, y, w, h) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLORS["white"]
        shape.line.color.rgb = COLORS["line"]
        band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, Inches(0.12))
        band.fill.solid()
        band.fill.fore_color.rgb = self.color(accent)
        band.line.fill.background()
        title_box = slide.shapes.add_textbox(x + Inches(0.18), y + Inches(0.15), w - Inches(0.3), Inches(0.35))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = FONT
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = COLORS["navy"]
        self.add_bullets(slide, bullets, x + Inches(0.18), y + Inches(0.55), w - Inches(0.35), h - Inches(0.7), size=12)

    def add_table(self, slide, headers: list[str], rows: list[list[str]], x, y, w, h, *, accent: str = "navy") -> None:
        table = slide.shapes.add_table(len(rows) + 1, len(headers), x, y, w, h).table
        for c, header in enumerate(headers):
            cell = table.cell(0, c)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.color(accent)
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.name = FONT
                    r.font.size = Pt(11)
                    r.font.bold = True
                    r.font.color.rgb = COLORS["white"]
        for r_idx, row in enumerate(rows, start=1):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["white"] if r_idx % 2 else COLORS["bg"]
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.name = FONT
                        run.font.size = Pt(10.5)
                        run.font.color.rgb = COLORS["ink"]

    def add_bar_chart(self, slide, categories: list[str], values: list[float], x, y, w, h) -> None:
        data = CategoryChartData()
        data.categories = categories
        data.add_series("Value", tuple(values))
        chart = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, x, y, w, h, data).chart
        chart.has_legend = False
        chart.value_axis.maximum_scale = 100
        chart.value_axis.minimum_scale = 0
        chart.value_axis.tick_labels.number_format = '0"%"'
        chart.category_axis.tick_labels.font.size = Pt(11)
        chart.value_axis.tick_labels.font.size = Pt(10)
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = COLORS["teal"]

    def render(self) -> None:
        case = self.case
        dp = self.spec["decision_pivot"]
        ctx = self.spec["company_industry_context"]
        dilemma = self.spec["core_dilemma"]
        evidence = self.spec["evidence"]
        options = self.spec["options"]
        appendix = self.spec.get("appendix", {})
        recommendation = self.spec.get("recommendation", {})

        slide = self.prs.slides.add_slide(self.layout)
        self.add_bg(slide, "bg")
        banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, Inches(1.25))
        banner.fill.solid()
        banner.fill.fore_color.rgb = COLORS["navy"]
        banner.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.32), Inches(11.8), Inches(0.55))
        p = title_box.text_frame.paragraphs[0]
        r = p.add_run()
        r.text = case["title"]
        r.font.name = FONT
        r.font.size = Pt(28)
        r.font.bold = True
        r.font.color.rgb = COLORS["white"]
        subtitle_lines = [dp["decision_question"]]
        if case.get("subtitle"):
            subtitle_lines.append(case["subtitle"])
        sub = slide.shapes.add_textbox(Inches(0.62), Inches(1.55), Inches(10.8), Inches(1.2))
        tf = sub.text_frame
        for idx, line in enumerate(subtitle_lines):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            r = p.add_run()
            r.text = line
            r.font.name = FONT
            r.font.size = Pt(18 if idx == 0 else 12)
            r.font.color.rgb = COLORS["ink"] if idx == 0 else COLORS["muted"]
        hero_cards = [
            ("決策時點", [dp["decision_owner"], dp["urgency"]], "blue"),
            ("真正張力", dilemma.get("trade_offs", [])[:2] or [dilemma["key_tension"]], "teal"),
            ("建議主張", [recommendation.get("recommended_option", "待決定"), recommendation.get("reason", "先整理 options 再表態")], "gold"),
        ]
        xs = [Inches(0.62), Inches(4.45), Inches(8.28)]
        for (title, bullets, accent), x in zip(hero_cards, xs):
            self.add_card(slide, title, bullets, accent, x, Inches(3.0), Inches(3.2), Inches(2.25))
        self.add_footer(slide, "Generated from canonical case spec")

        slide = self.prs.slides.add_slide(self.layout)
        self.add_bg(slide, "bg")
        self.add_header(slide, "1. 導言與決策切入點", "主角、時間軸與延遲代價")
        self.add_card(slide, "決策主角", [dp["decision_owner"], dp["decision_question"]], "blue", Inches(0.55), Inches(1.45), Inches(4.0), Inches(4.95))
        self.add_card(slide, "時間軸", dp.get("milestones", []), "teal", Inches(4.75), Inches(1.45), Inches(4.0), Inches(4.95))
        self.add_card(slide, "延遲代價", [dp["urgency"], dp["delay_cost"]], "red", Inches(8.95), Inches(1.45), Inches(3.85), Inches(4.95))
        self.add_footer(slide)

        slide = self.prs.slides.add_slide(self.layout)
        self.add_bg(slide, "white")
        self.add_header(slide, "2. 公司與產業總體檢", "公司基礎、產業背景與結構拆解")
        cards = [
            ("公司基礎", ctx.get("company_background", []), "blue"),
            ("產業背景", ctx.get("industry_background", []), "teal"),
            ("結構拆解", sum(ctx.get("structural_breakdown", {}).values(), []), "gold"),
        ]
        xs = [Inches(0.55), Inches(4.75), Inches(8.95)]
        for (title, bullets, accent), x in zip(cards, xs):
            self.add_card(slide, title, bullets[:4], accent, x, Inches(1.45), Inches(3.85), Inches(4.85))
        self.add_footer(slide)

        slide = self.prs.slides.add_slide(self.layout)
        self.add_bg(slide, "bg")
        self.add_header(slide, "3. 核心衝突與根本問題", "表層問題、根因與 trade-off")
        self.add_table(
            slide,
            ["面向", "內容"],
            [
                ["表層問題", dilemma["surface_problem"]],
                ["根本問題", dilemma["root_problem"]],
                ["核心張力", dilemma["key_tension"]],
            ],
            Inches(0.65), Inches(1.55), Inches(7.1), Inches(2.2),
        )
        self.add_bullets(slide, ["5 Whys"] + dilemma.get("five_whys", []), Inches(8.0), Inches(1.6), Inches(4.3), Inches(2.8), size=13)
        self.add_card(slide, "Trade-off", dilemma.get("trade_offs", []), "red", Inches(0.75), Inches(4.3), Inches(11.6), Inches(1.9))
        self.add_footer(slide)

        slide = self.prs.slides.add_slide(self.layout)
        self.add_bg(slide, "white")
        self.add_header(slide, "4. 數據驗證與實例支持", "關鍵數據、內部檢查與外部交叉")
        if evidence.get("quantitative_signals"):
            values = [max(10, 100 - idx * 15) for idx, _ in enumerate(evidence["quantitative_signals"][:5])]
            labels = [f"S{idx+1}" for idx, _ in enumerate(values)]
            self.add_bar_chart(slide, labels, values, Inches(7.0), Inches(1.6), Inches(5.4), Inches(3.2))
        self.add_bullets(slide, evidence.get("quantitative_signals", []), Inches(0.7), Inches(1.6), Inches(5.8), Inches(2.7), size=14)
        self.add_card(slide, "內部檢查", evidence.get("internal_checks", []), "gold", Inches(0.75), Inches(4.8), Inches(5.7), Inches(1.7))
        self.add_card(slide, "外部交叉", evidence.get("external_checks", []), "blue", Inches(6.85), Inches(4.8), Inches(5.6), Inches(1.7))
        self.add_footer(slide)

        slide = self.prs.slides.add_slide(self.layout)
        self.add_bg(slide, "white")
        self.add_header(slide, "5. 課堂思辨與行動決策", "Options 與 recommendation")
        cols = min(3, len(options))
        gap = Inches(0.2)
        card_w = (Inches(12.2) - gap * (cols - 1)) / cols
        for idx, option in enumerate(options[:3]):
            bullets = [f"How：{option['how']}", f"Why：{option['why']}", f"Trade-off：{option['trade_off']}"]
            self.add_card(slide, option["title"], bullets, ["teal", "gold", "blue"][idx % 3], Inches(0.55) + idx * (card_w + gap), Inches(1.45), card_w, Inches(4.8))
        self.add_footer(slide)

        slide = self.prs.slides.add_slide(self.layout)
        self.add_bg(slide, "bg")
        self.add_header(slide, "6. 建議與附錄", "推薦方案、前提與討論題")
        self.add_card(slide, "推薦方案", [
            f"建議：{recommendation.get('recommended_option', '未指定')}",
            f"理由：{recommendation.get('reason', '未指定')}",
        ] + recommendation.get("conditions", []), "teal", Inches(0.6), Inches(1.45), Inches(4.0), Inches(4.8))
        self.add_card(slide, "參考與假設", appendix.get("references", [])[:3] + appendix.get("assumptions", [])[:2], "gold", Inches(4.8), Inches(1.45), Inches(3.9), Inches(4.8))
        self.add_card(slide, "課堂討論題", appendix.get("discussion_questions", [])[:4] + self.case.get("source_notes", [])[:2], "blue", Inches(8.9), Inches(1.45), Inches(3.9), Inches(4.8))
        self.add_footer(slide)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Generate PPTX from the canonical case spec.")
    parser.add_argument("spec", type=Path, help="Path to canonical case spec JSON.")
    parser.add_argument("output", type=Path, help="Path to output .pptx file.")
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    spec = load_spec(args.spec)
    builder = DeckBuilder(spec)
    builder.render()
    output_path = args.output.expanduser().resolve()
    output_path.parent.mkdir(parents=True, exist_ok=True)
    builder.prs.save(output_path)
    print(output_path)


if __name__ == "__main__":
    main()
